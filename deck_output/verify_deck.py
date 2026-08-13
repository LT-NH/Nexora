"""Verify the generated PPTX by dumping slide structure, shape positions, and text."""

from pptx import Presentation

prs = Presentation(r"C:/Users/lihaoqi/Desktop/SaaS/deck_output/Nexora-Product-Launch.pptx")
print(f"Slide size: {prs.slide_width} x {prs.slide_height} EMU")
print(f"Slide size: {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} inches")
print(f"Total slides: {len(prs.slides)}")
print("=" * 80)

for i, slide in enumerate(prs.slides, 1):
    print(f"\n[Slide {i}]  shape count = {len(slide.shapes)}")
    for j, shp in enumerate(slide.shapes):
        kind = shp.shape_type
        x_in = shp.left / 914400 if shp.left else 0
        y_in = shp.top / 914400 if shp.top else 0
        w_in = shp.width / 914400 if shp.width else 0
        h_in = shp.height / 914400 if shp.height else 0
        txt = ""
        if shp.has_text_frame:
            txt = " | ".join(p.text for p in shp.text_frame.paragraphs)[:80]
        print(f"  #{j:02d} {str(kind):40s} @({x_in:5.2f},{y_in:5.2f}) {w_in:5.2f}x{h_in:5.2f}  '{txt}'")