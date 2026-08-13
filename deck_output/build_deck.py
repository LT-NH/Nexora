"""Build the Nexora product launch deck from the 7 AI-generated background images.

Strategy:
- Use each AI image as a full-bleed background on the slide.
- Add a dark gradient overlay rectangle on the left half to ensure text legibility.
- Place a solid dark patch over the AI watermark at the bottom-right corner.
- Overlay Chinese titles, body copy, accent lines, and brand badge using python-pptx.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

OUT_DIR = Path(r"C:/Users/lihaoqi/Desktop/SaaS/deck_output")
IMG_DIR = OUT_DIR / "images"
OUTPUT = OUT_DIR / "Nexora-Product-Launch.pptx"

# Brand palette tuned to the AI-generated dark-tech imagery
BG_NAVY = RGBColor(0x0B, 0x10, 0x1F)        # deep base
PANEL = RGBColor(0x12, 0x1A, 0x2E)          # glass panel
ACCENT_CYAN = RGBColor(0x22, 0xE6, 0xFF)    # neon cyan
ACCENT_VIOLET = RGBColor(0xA0, 0x6B, 0xFF)  # violet
ACCENT_GOLD = RGBColor(0xFF, 0xC4, 0x5C)    # highlight gold
TEXT_PRIMARY = RGBColor(0xF5, 0xF8, 0xFF)
TEXT_SECONDARY = RGBColor(0x9D, 0xA8, 0xC2)
TEXT_MUTED = RGBColor(0x6E, 0x78, 0x96)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Slide content
SLIDES = [
    {
        "image": "01-cover.png",
        "kind": "cover",
        "eyebrow": "PRODUCT LAUNCH · 2026",
        "title": "Nexora",
        "subtitle": "One Core, All Commerce",
        "desc": "面向中小电商的一站式多租户 SaaS 管理平台",
        "footer": "李浩棋 · 全栈独立开发 ·  github.com/lihaoqi-13656117061/nexora",
        "accent": ACCENT_CYAN,
    },
    {
        "image": "02-painpoints.png",
        "kind": "content",
        "eyebrow": "01 · 痛点",
        "title": "中小商家的电商之痛",
        "points": [
            ("多平台割裂", "Shopify、抖音、淘宝…切换多个后台，\n数据零散难统一"),
            ("数据孤岛", "商品、订单、客户散落各处，\n无法跨平台汇总分析"),
            ("管理低效", "重复录入、人工对账、缺乏洞察，\n被运营琐事淹没"),
        ],
        "footer": "中小商家需要一个真正统一的后台",
        "accent": ACCENT_CYAN,
    },
    {
        "image": "03-solution.png",
        "kind": "content",
        "eyebrow": "02 · 解决方案",
        "title": "Nexora — 一个后台，统管全域电商",
        "points": [
            ("统一中枢", "多租户工作空间，一个账号聚合所有店铺的\n商品、订单、客户、AI 洞察"),
            ("全链路覆盖", "从商品上架、订单流转、客户 CRM 到数据\n分析与 AI 赋能，端到端打通"),
            ("开箱即用", "Docker 一键部署，演示数据秒级生成\n完整业务场景"),
        ],
        "footer": "不再切来切去 —— 一个后台就是全部",
        "accent": ACCENT_VIOLET,
    },
    {
        "image": "04-features.png",
        "kind": "features",
        "eyebrow": "03 · 核心功能",
        "title": "六大模块，一个工作台",
        "points": [
            ("商品管理", "SKU / 变体 / 分类树 · AI 描述生成"),
            ("订单管理", "全状态流转 · 跨平台趋势统计"),
            ("客户 CRM", "标签 + RFM 五层价值分析"),
            ("多平台集成", "Shopify / 抖音 / 淘宝 / 京东 / Amazon"),
            ("AI 洞察", "销售预测 · 营销文案 · SEO 关键词"),
            ("企业级安全", "JWT + Fernet + bcrypt + 速率限制"),
        ],
        "footer": "20 页面 · 17 UI 组件 · 6 类 ECharts 图表",
        "accent": ACCENT_CYAN,
    },
    {
        "image": "05-advantage.png",
        "kind": "advantages",
        "eyebrow": "04 · 竞争优势",
        "title": "轻量 · 全域 · AI 三位一体",
        "points": [
            ("vs 单平台工具", "覆盖更广，一个后台管 Shopify + 抖音\n+ 淘宝 + 京东 + Amazon"),
            ("vs 大型 ERP", "更轻量、更易上手，5 分钟完成\nDocker 一键部署"),
            ("核心壁垒", "AI 能力深度嵌入业务场景，\n让每个商家拥有智能助手"),
        ],
        "footer": "聚焦中小商家 (年 GMV 50 万 — 5000 万)",
        "accent": ACCENT_GOLD,
    },
    {
        "image": "06-pricing.png",
        "kind": "pricing",
        "eyebrow": "05 · 商业模式 / 定价",
        "title": "订阅制 SaaS，三档套餐成长陪伴",
        "tiers": [
            ("Free", "¥0", "/月", ["基础商品 / 订单管理", "1 个工作空间", "社区支持"]),
            ("Pro", "¥199", "/月", ["全模块功能解锁", "5 个工作空间", "AI 洞察 + 多平台集成", "工单优先响应"]),
            ("Enterprise", "¥699", "/月", ["无限工作空间", "专属客户经理", "API / 插件扩展", "SLA 99.9%"]),
        ],
        "footer": "低门槛获客 → 数据沉淀 → 套餐升级 · 增长飞轮",
        "accent": ACCENT_CYAN,
    },
    {
        "image": "07-cta.png",
        "kind": "cta",
        "eyebrow": "06 · 现在就开始",
        "title": "欢迎共建下一代电商中台",
        "subtitle": "Demo · 合作 · 投资",
        "points": [
            ("在线体验", "demo.nexora.com  ·  demo@nexora.com / Demo1234!"),
            ("开源仓库", "github.com/lihaoqi-13656117061/nexora"),
            ("一键部署", "bash deploy.sh  →  localhost:3000"),
            ("商务联系", "李浩棋 ·  13656117061"),
        ],
        "footer": "感谢聆听 · 期待与您共创",
        "accent": ACCENT_CYAN,
    },
]


def add_rect(slide, x, y, w, h, fill, line=None, transparency=None):
    """Add a solid-fill rectangle (optionally with transparency on fill)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.shadow.inherit = False
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if transparency is not None:
        # Apply transparency via alpha on solid fill
        sp = shape.fill.fore_color._xFill
        srgb = sp.getparent().find("{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
        if srgb is not None:
            color_el = srgb.find("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
            if color_el is not None:
                # Insert alpha child element
                from lxml import etree
                alpha = etree.SubElement(
                    color_el,
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha",
                )
                alpha.set("val", str(int(transparency * 1000)))
    return shape


def set_text(
    text_frame,
    text,
    *,
    font="Microsoft YaHei",
    size=18,
    color=TEXT_PRIMARY,
    bold=False,
    align=None,
):
    """Replace text frame content with a single paragraph."""
    text_frame.clear()
    p = text_frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Disable Asian-font auto-substitution quirks by also setting east-asian font
    rpr = run._r.get_or_add_rPr()
    from lxml import etree
    ea = rpr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
    )
    if ea is None:
        ea = etree.SubElement(
            rpr,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
        )
    ea.set("typeface", font)
    return run


def add_textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    tb.text_frame.margin_left = Emu(0)
    tb.text_frame.margin_right = Emu(0)
    tb.text_frame.margin_top = Emu(0)
    tb.text_frame.margin_bottom = Emu(0)
    return tb


def add_run(text_frame, text, *, font="Microsoft YaHei", size=18, color=TEXT_PRIMARY, bold=False):
    """Append a styled run to the text frame, creating a new paragraph if needed."""
    if text_frame.paragraphs[0].runs and not text_frame.paragraphs[0].text:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    from lxml import etree
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
    )
    if ea is None:
        ea = etree.SubElement(
            rpr,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
        )
    ea.set("typeface", font)
    return run


def layout_slide(prs, slide_def):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    width = prs.slide_width
    height = prs.slide_height

    # 1. Background image full-bleed
    slide.shapes.add_picture(
        str(IMG_DIR / slide_def["image"]),
        0,
        0,
        width=width,
        height=height,
    )

    # 2. Bottom-right watermark mask (covers "图片由AI生成")
    mask_w = Inches(2.4)
    mask_h = Inches(0.5)
    mask_x = width - mask_w - Inches(0.05)
    mask_y = height - mask_h - Inches(0.05)
    add_rect(slide, mask_x, mask_y, mask_w, mask_h, fill=BG_NAVY)

    # 3. Subtle dark gradient strip on the left half (panel feel)
    panel_w = int(width * 0.58)
    add_rect(slide, 0, 0, panel_w, height, fill=BG_NAVY, transparency=0.32)

    # 4. Brand badge top-left
    badge = add_textbox(slide, Inches(0.7), Inches(0.55), Inches(3), Inches(0.4))
    add_run(
        badge.text_frame,
        "◆ NEXORA",
        size=14,
        color=slide_def["accent"],
        bold=True,
    )

    # 5. Eyebrow
    eyebrow_tb = add_textbox(slide, Inches(0.7), Inches(1.15), Inches(6), Inches(0.4))
    add_run(
        eyebrow_tb.text_frame,
        slide_def["eyebrow"],
        size=12,
        color=TEXT_MUTED,
        bold=True,
    )

    # 6. Title
    title_tb = add_textbox(slide, Inches(0.7), Inches(1.6), int(panel_w * 0.92), Inches(1.4))
    add_run(
        title_tb.text_frame,
        slide_def["title"],
        size=40,
        color=TEXT_PRIMARY,
        bold=True,
    )

    # 7. Accent underline
    accent_line = add_rect(
        slide,
        Inches(0.7),
        Inches(2.95),
        Inches(0.9),
        Inches(0.06),
        fill=slide_def["accent"],
    )

    kind = slide_def["kind"]
    base_y = Inches(3.25)

    if kind == "cover":
        # Big subtitle, desc, and footer
        sub_tb = add_textbox(slide, Inches(0.7), base_y, int(panel_w * 0.9), Inches(1.0))
        add_run(
            sub_tb.text_frame,
            slide_def["subtitle"],
            size=30,
            color=slide_def["accent"],
            bold=True,
        )
        desc_tb = add_textbox(slide, Inches(0.7), base_y + Inches(1.05), int(panel_w * 0.9), Inches(0.8))
        add_run(desc_tb.text_frame, slide_def["desc"], size=18, color=TEXT_SECONDARY)
        footer_tb = add_textbox(slide, Inches(0.7), height - Inches(0.7), int(panel_w * 0.9), Inches(0.4))
        add_run(footer_tb.text_frame, slide_def["footer"], size=11, color=TEXT_MUTED)
        # Right-side stack of key facts
        fact_x = int(panel_w + Inches(0.5))
        fact_w = width - fact_x - Inches(0.5)
        facts = [
            ("147", "源码文件"),
            ("28/28", "测试通过"),
            ("12", "API 路由"),
            ("20", "前端页面"),
        ]
        for idx, (num, label) in enumerate(facts):
            ftb = add_textbox(slide, fact_x, Inches(1.6 + idx * 1.05), fact_w, Inches(1.0))
            tf = ftb.text_frame
            add_run(tf, num, size=32, color=slide_def["accent"], bold=True)
            add_run(tf, label, size=12, color=TEXT_SECONDARY)
        return

    if kind == "content":
        for idx, (head, body) in enumerate(slide_def["points"]):
            tb = add_textbox(
                slide,
                Inches(0.7),
                base_y + Inches(idx * 1.25),
                int(panel_w * 0.9),
                Inches(1.1),
            )
            tf = tb.text_frame
            add_run(tf, "▍ " + head, size=18, color=slide_def["accent"], bold=True)
            add_run(tf, body, size=14, color=TEXT_PRIMARY)
        footer_tb = add_textbox(
            slide,
            Inches(0.7),
            height - Inches(0.7),
            int(panel_w * 0.9),
            Inches(0.4),
        )
        add_run(footer_tb.text_frame, slide_def["footer"], size=11, color=TEXT_MUTED)
        return

    if kind == "features":
        # 2 columns x 3 rows of compact feature pills
        col_w = (panel_w - Inches(0.9) - Inches(0.7)) / 2
        row_h = Inches(1.05)
        for idx, (head, body) in enumerate(slide_def["points"]):
            r = idx // 2
            c = idx % 2
            x = Inches(0.7) + c * (col_w + Inches(0.3))
            y = base_y + r * row_h
            # Pill background
            pill = add_rect(
                slide,
                x,
                y,
                int(col_w),
                int(row_h) - Emu(80000),
                fill=PANEL,
                transparency=0.45,
            )
            pill.line.color.rgb = slide_def["accent"]
            pill.line.width = Pt(0.5)
            pill.line.fill.solid()
            pill.line.fill.fore_color.rgb = slide_def["accent"]
            tb = add_textbox(
                slide,
                x + Inches(0.15),
                y + Inches(0.1),
                int(col_w) - Inches(0.3),
                int(row_h) - Inches(0.2),
            )
            tf = tb.text_frame
            add_run(tf, head, size=14, color=slide_def["accent"], bold=True)
            add_run(tf, body, size=11, color=TEXT_PRIMARY)
        footer_tb = add_textbox(
            slide,
            Inches(0.7),
            height - Inches(0.7),
            int(panel_w * 0.9),
            Inches(0.4),
        )
        add_run(footer_tb.text_frame, slide_def["footer"], size=11, color=TEXT_MUTED)
        return

    if kind == "advantages":
        for idx, (head, body) in enumerate(slide_def["points"]):
            tb = add_textbox(
                slide,
                Inches(0.7),
                base_y + Inches(idx * 1.25),
                int(panel_w * 0.9),
                Inches(1.1),
            )
            tf = tb.text_frame
            add_run(tf, "◆ " + head, size=18, color=slide_def["accent"], bold=True)
            add_run(tf, body, size=14, color=TEXT_PRIMARY)
        footer_tb = add_textbox(
            slide,
            Inches(0.7),
            height - Inches(0.7),
            int(panel_w * 0.9),
            Inches(0.4),
        )
        add_run(footer_tb.text_frame, slide_def["footer"], size=11, color=TEXT_MUTED)
        return

    if kind == "pricing":
        tier_x = Inches(0.7)
        tier_y = base_y
        tier_w = (panel_w - Inches(1.4)) / 3
        tier_h = Inches(3.2)
        tier_accent = [ACCENT_CYAN, ACCENT_VIOLET, ACCENT_GOLD]
        for idx, (name, price, unit, features) in enumerate(slide_def["tiers"]):
            x = tier_x + idx * (tier_w + Inches(0.3))
            card = add_rect(
                slide,
                int(x),
                int(tier_y),
                int(tier_w),
                int(tier_h),
                fill=PANEL,
                transparency=0.5,
            )
            card.line.color.rgb = tier_accent[idx]
            card.line.width = Pt(1.0)
            card.line.fill.solid()
            card.line.fill.fore_color.rgb = tier_accent[idx]

            # top accent bar
            bar_h = Inches(0.12)
            add_rect(
                slide,
                int(x),
                int(tier_y),
                int(tier_w),
                int(bar_h),
                fill=tier_accent[idx],
            )

            tb = add_textbox(
                slide,
                int(x) + Inches(0.25),
                int(tier_y) + Inches(0.3),
                int(tier_w) - Inches(0.5),
                int(tier_h) - Inches(0.4),
            )
            tf = tb.text_frame
            add_run(tf, name, size=18, color=tier_accent[idx], bold=True)
            add_run(tf, price, size=28, color=TEXT_PRIMARY, bold=True)
            add_run(tf, unit, size=12, color=TEXT_MUTED)
            for f in features:
                add_run(tf, "✓ " + f, size=11, color=TEXT_PRIMARY)
        footer_tb = add_textbox(
            slide,
            Inches(0.7),
            height - Inches(0.7),
            int(panel_w * 0.9),
            Inches(0.4),
        )
        add_run(footer_tb.text_frame, slide_def["footer"], size=11, color=TEXT_MUTED)
        return

    if kind == "cta":
        sub_tb = add_textbox(slide, Inches(0.7), base_y, int(panel_w * 0.9), Inches(0.7))
        add_run(sub_tb.text_frame, slide_def["subtitle"], size=24, color=slide_def["accent"], bold=True)
        for idx, (head, body) in enumerate(slide_def["points"]):
            tb = add_textbox(
                slide,
                Inches(0.7),
                base_y + Inches(0.8) + Inches(idx * 0.7),
                int(panel_w * 0.92),
                Inches(0.65),
            )
            tf = tb.text_frame
            add_run(tf, "▍ " + head, size=14, color=slide_def["accent"], bold=True)
            add_run(tf, body, size=12, color=TEXT_PRIMARY)
        footer_tb = add_textbox(
            slide,
            Inches(0.7),
            height - Inches(0.45),
            int(panel_w * 0.9),
            Inches(0.35),
        )
        add_run(footer_tb.text_frame, slide_def["footer"], size=11, color=TEXT_MUTED)
        return


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for s in SLIDES:
        layout_slide(prs, s)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()